"""
FLEETSYNC TMS — Sunum Oluşturucu
python-pptx ile detaylı kullanım kılavuzu sunumu
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Renk paleti (FLEETSYNC dark tema) ───────────────────────────
BG       = RGBColor(0x0a, 0x0b, 0x0a)   # #0a0b0a
PANEL    = RGBColor(0x15, 0x17, 0x1a)   # #15171a
LINE     = RGBColor(0x2a, 0x2e, 0x33)   # #2a2e33
AMBER    = RGBColor(0xF5, 0xB5, 0x00)   # amber
GREEN    = RGBColor(0x4C, 0xD9, 0x6C)   # green
RED      = RGBColor(0xE5, 0x4D, 0x42)   # red
INK      = RGBColor(0xe8, 0xe6, 0xe0)   # #e8e6e0
INK_DIM  = RGBColor(0xa0, 0xa3, 0x9a)   # #a0a39a
INK_MUTE = RGBColor(0x6b, 0x6f, 0x6a)   # #6b6f6a
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank


# ── Yardımcı fonksiyonlar ────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line_color=None, line_width=Pt(0.5)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.color.rgb = line_color or LINE
    shape.line.width = line_width
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=11, bold=False, color=INK, align=PP_ALIGN.LEFT,
             font_name="Courier New", wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txb

def add_tag(slide, text, l, t, color=AMBER, size=8):
    w = len(text) * 0.075 + 0.15
    box = add_rect(slide, l, t, w, 0.22, fill=None, line_color=color)
    add_text(slide, text, l + 0.04, t + 0.01, w - 0.08, 0.2,
             size=size, color=color, align=PP_ALIGN.CENTER)
    return w

def slide_bg(slide):
    add_rect(slide, 0, 0, 13.33, 7.5, fill=BG, line_color=BG)

def topbar(slide, active_tab="Overview"):
    # top strip
    add_rect(slide, 0, 0, 13.33, 0.38, fill=PANEL, line_color=LINE)
    add_text(slide, "FLEETSYNC", 0.2, 0.04, 1.6, 0.3,
             size=16, bold=True, color=AMBER, font_name="Arial Black")
    add_text(slide, "/DISPATCH·V0.2", 1.82, 0.11, 1.5, 0.2,
             size=7, color=INK_MUTE)
    add_text(slide, "● SYS ONLINE", 3.4, 0.1, 1.2, 0.2, size=8, color=AMBER)
    add_text(slide, "● API · 200 OK", 4.65, 0.1, 1.2, 0.2, size=8, color=GREEN)
    add_text(slide, "● BOT · ACTIVE", 5.9, 0.1, 1.2, 0.2, size=8, color=GREEN)
    add_text(slide, "TIR FLEET  |  M YILMAZ  |  OWNER  |  LOGOUT",
             9.5, 0.1, 3.7, 0.2, size=8, color=INK_MUTE, align=PP_ALIGN.RIGHT)

    # nav tabs
    add_rect(slide, 0, 0.38, 13.33, 0.38, fill=BG, line_color=LINE)
    tabs = [("F1","Overview"),("F2","PTI Inspection"),("F3","Loads / Dispatch"),
            ("F4","Fleet Roster"),("F5","Factoring"),("F6","IFTA")]
    x = 0.0
    for key, label in tabs:
        is_active = (label == active_tab)
        w = len(label) * 0.095 + 0.5
        if is_active:
            add_rect(slide, x, 0.38, w, 0.38, fill=PANEL, line_color=AMBER,
                     line_width=Pt(1.5))
        add_text(slide, key, x + 0.1, 0.44, 0.2, 0.24, size=7, color=INK_MUTE)
        add_text(slide, label, x + 0.32, 0.43, w - 0.4, 0.24, size=9,
                 color=AMBER if is_active else INK)
        x += w

def panel_head(slide, text, l, t, w, right_text=""):
    add_rect(slide, l, t, w, 0.28, fill=PANEL, line_color=LINE)
    add_text(slide, "· " + text, l + 0.1, t + 0.04, w * 0.6, 0.22,
             size=8, color=INK_DIM)
    if right_text:
        add_text(slide, right_text, l + w * 0.55, t + 0.04, w * 0.42, 0.22,
                 size=8, color=INK_MUTE, align=PP_ALIGN.RIGHT)

def callout(slide, text, l, t, w=3.5, color=AMBER):
    """Açıklama kutusu — sarı kenarlıklı"""
    add_rect(slide, l, t, w, 0.32, fill=RGBColor(0x1b,0x1e,0x22), line_color=color)
    add_text(slide, "▸ " + text, l + 0.1, t + 0.05, w - 0.2, 0.24,
             size=8.5, color=color)

def numbered_list(slide, items, l, t, w=4.0, color=INK_DIM):
    for i, (title, desc) in enumerate(items):
        y = t + i * 0.42
        add_text(slide, f"{i+1:02d}", l, y, 0.28, 0.35, size=9, color=AMBER, bold=True)
        add_text(slide, title, l + 0.3, y, w - 0.3, 0.18, size=9, bold=True, color=INK)
        add_text(slide, desc, l + 0.3, y + 0.18, w - 0.3, 0.22, size=7.5, color=INK_DIM)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — Kapak
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)

# Sol dikey şerit
add_rect(s, 0, 0, 0.06, 7.5, fill=AMBER, line_color=AMBER)

# Logo büyük
add_text(s, "FLEETSYNC", 0.4, 1.6, 8, 1.4,
         size=72, bold=True, color=AMBER, font_name="Arial Black")
add_text(s, "/DISPATCH · V0.2", 0.45, 2.95, 6, 0.4,
         size=14, color=INK_MUTE, font_name="Courier New")

# Tagline
add_text(s, "Transportation Management System", 0.45, 3.5, 8, 0.5,
         size=18, color=INK, font_name="Arial")
add_text(s, "Filo Yönetimi · PTI · Dispatch · Factoring · DAT Load Board · IFTA",
         0.45, 4.0, 9, 0.4, size=11, color=INK_DIM)

# Feature tags
tags = ["DISPATCH","PTI","DAT BOARD","FACTORING","DOCS","MAINTENANCE","IFTA","MULTI-TENANT"]
x = 0.45
for tag in tags:
    tw = add_tag(s, tag, x, 4.6, color=AMBER, size=8)
    x += tw + 0.15

# Sağ panel — sistem bilgisi
add_rect(s, 9.5, 1.4, 3.6, 4.8, fill=PANEL, line_color=LINE)
add_text(s, "· SİSTEM BİLGİSİ", 9.65, 1.55, 3.3, 0.25, size=8, color=INK_MUTE)
stats = [
    ("Stack",       "FastAPI + React + SQLite"),
    ("Bot",         "Telegram (PTI)"),
    ("Server",      "Hetzner VPS · 5.161.61.221"),
    ("Auth",        "JWT Bearer Token"),
    ("Mimari",      "Multi-Tenant (Company)"),
    ("Versiyon",    "0.2 · Mayıs 2026"),
]
for i, (k, v) in enumerate(stats):
    y = 1.9 + i * 0.55
    add_text(s, k, 9.65, y, 1.1, 0.25, size=8, color=INK_MUTE)
    add_text(s, v, 10.8, y, 2.2, 0.25, size=8, color=INK)
    if i < len(stats)-1:
        add_rect(s, 9.65, y + 0.32, 3.3, 0.01, fill=LINE, line_color=LINE)

# Alt çizgi
add_rect(s, 0.4, 6.8, 12.9, 0.01, fill=LINE, line_color=LINE)
add_text(s, "GİZLİ VE ÖZEL · YETKİSİZ DAĞITIM YASAKTIR · 2026",
         0.4, 6.9, 12, 0.3, size=8, color=INK_MUTE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — İçindekiler / Modüller
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
add_rect(s, 0, 0, 0.06, 7.5, fill=AMBER, line_color=AMBER)

add_text(s, "FLEETSYNC — MODÜL HARİTASI", 0.3, 0.2, 12, 0.4,
         size=18, bold=True, color=AMBER, font_name="Arial Black")
add_text(s, "Sistemin tüm modülleri ve kapsamı",
         0.3, 0.62, 10, 0.3, size=10, color=INK_DIM)
add_rect(s, 0.3, 0.95, 12.7, 0.01, fill=LINE, line_color=LINE)

modules = [
    ("F1", "OVERVİEW / DASHBOARD", GREEN,
     "Filo durumu anlık görünümü. Tüm tırlar, PTI durumu, aktif yükler ve gelir özeti.\nHer 30 saniyede otomatik güncellenir."),
    ("F2", "PTI INSPECTION", AMBER,
     "Günlük pre-trip inspection takibi. Telegram botu üzerinden fotoğraf alınır.\nKim gönderdi, kaç fotoğraf, hangi saatte bilgisi görüntülenir."),
    ("F3", "LOADS / DISPATCH", AMBER,
     "Yük oluşturma, durum takibi, PDF rate conf parse, DAT board entegrasyonu.\nBOL ve belge yükleme, yük arama filtresi."),
    ("F4", "FLEET ROSTER", GREEN,
     "Tır ve şoför yönetimi. Kilometre, yakıt, VIN, CDL expiry takibi.\nBakım log sistemi — servis geçmişi, maliyet kaydı."),
    ("F5", "FACTORING", AMBER,
     "RTS factoring iş akışı. Teslim edilen yükten fatura oluşturma, %3.5 komisyon hesabı.\nSubmit → Approved → Paid takip ekranı."),
    ("F6", "IFTA", INK_DIM,
     "Quarter bazlı state mileage özeti. Yakıt vergisi tahmini.\nIFTA quarterly filing için ön rapor."),
]

for i, (key, title, color, desc) in enumerate(modules):
    col = i % 2
    row = i // 2
    l = 0.3 + col * 6.5
    t = 1.1 + row * 1.95

    add_rect(s, l, t, 6.2, 1.8, fill=PANEL, line_color=LINE)
    add_tag(s, key, l + 0.15, t + 0.15, color=color)
    add_text(s, title, l + 0.7, t + 0.12, 5.3, 0.3, size=11, bold=True, color=color,
             font_name="Arial Black")
    add_text(s, desc, l + 0.15, t + 0.5, 5.9, 1.1, size=9, color=INK_DIM)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — Dashboard (Overview)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
topbar(s, "Overview")

# Stats bar
stats_bar = [
    ("Active Trucks", "01", "/ 1 total", AMBER),
    ("PTI Today", "1/1", "all clear", GREEN),
    ("In Transit", "00", "rolling", INK),
    ("Gross · WTD", "$0.2K", "$0.04/mi", GREEN),
    ("On-Time", "50%", "last 30 days", AMBER),
    ("Open Alerts", "00", "none", INK),
]
for i, (label, val, sub, color) in enumerate(stats_bar):
    x = 0.0 + i * (13.33/6)
    w = 13.33/6
    add_rect(s, x, 0.76, w, 0.9, fill=PANEL, line_color=LINE)
    add_text(s, label, x + 0.1, 0.8, w - 0.15, 0.18, size=7, color=INK_MUTE)
    add_text(s, sub, x + 0.1, 0.8, w - 0.15, 0.18, size=7, color=INK_MUTE,
             align=PP_ALIGN.RIGHT)
    add_text(s, val, x + 0.1, 1.0, w - 0.2, 0.5, size=22, bold=True, color=color,
             font_name="Arial Black")

# Fleet status board
panel_head(s, "FLEET STATUS BOARD · LIVE", 0, 1.66, 9.0, "AUTO-REFRESH 30s")
add_rect(s, 0, 1.94, 9.0, 0.28, fill=PANEL, line_color=LINE)
for j, h in enumerate(["","Unit","Driver","Status","Lane / Position","Rate","ETA","PTI","Fuel"]):
    add_text(s, h, 0.1 + j*1.0, 1.97, 0.95, 0.22, size=7, color=INK_MUTE)
# Mock row
add_rect(s, 0, 2.22, 9.0, 0.45, fill=BG, line_color=LINE)
row_vals = ["●","9852","M. Bayram","ROLLING","Dallas TX → Chicago IL","$2,400","18:30 CT","PTI·OK","███ 80%"]
colors_r  = [AMBER,INK,INK_DIM,AMBER,INK_DIM,GREEN,INK_DIM,GREEN,GREEN]
for j, (v, c) in enumerate(zip(row_vals, colors_r)):
    add_text(s, v, 0.1 + j*1.0, 2.28, 0.95, 0.32, size=8, color=c)

# Callout açıklamaları
callout(s, "Her tır için anlık durum — status, yük, PTI, yakıt", 0.15, 2.75, 5.5, AMBER)
callout(s, "Yeşil dot = ROLLING  |  Kırmızı = SHOP  |  Gri = IDLE", 0.15, 3.15, 5.5, GREEN)

# Alert feed
panel_head(s, "ALERT FEED", 9.1, 1.66, 4.1, "4 ITEMS")
alerts = [
    (RED,   "0 DRIVER(S) MISSING PTI TODAY"),
    (AMBER, "FUEL ALERT · TRUCK(S) BELOW 25%"),
    (INK_MUTE, "FLEET · 1 TOTAL UNITS REGISTERED"),
    (INK_MUTE, "RTS FACTORING · CHECK PENDING"),
]
for i, (c, txt) in enumerate(alerts):
    y = 1.94 + i * 0.35
    add_rect(s, 9.1, y, 4.1, 0.34, fill=BG, line_color=LINE)
    add_text(s, "●", 9.2, y + 0.08, 0.2, 0.2, size=7, color=c)
    add_text(s, txt, 9.45, y + 0.07, 3.65, 0.22, size=8, color=c)

# PTI bar chart
panel_head(s, "PTI COMPLIANCE · TODAY", 9.1, 3.3, 4.1, "100%")
add_rect(s, 9.1, 3.58, 4.1, 0.7, fill=BG, line_color=LINE)
add_rect(s, 9.5, 3.65, 0.35, 0.5, fill=GREEN, line_color=GREEN)
add_text(s, "MUST", 9.5, 4.18, 0.4, 0.15, size=7, color=INK_MUTE)

# Quick commands
panel_head(s, "QUICK COMMAND", 9.1, 4.35, 4.1)
cmds = ["▸ NUDGE LATE PTI", "+ NEW LOAD", "+ DISPATCH",
        "RTS BATCH", "DRIVER MSG", "EXPORT CSV"]
for i, c in enumerate(cmds):
    col = i % 2
    row = i // 2
    add_rect(s, 9.15 + col*2.05, 4.63 + row*0.42, 1.95, 0.35,
             fill=AMBER if i==0 else BG, line_color=AMBER if i==0 else LINE)
    add_text(s, c, 9.2 + col*2.05, 4.7 + row*0.42, 1.85, 0.24,
             size=8, color=BG if i==0 else INK, align=PP_ALIGN.CENTER)

# Console
add_rect(s, 9.1, 5.9, 4.1, 1.35, fill=PANEL, line_color=LINE)
add_text(s, "> status --fleet\n  1 ACTIVE · 0 OFFLINE\n> pti --today\n  1/1 OK · 0 MISS\n_",
         9.2, 5.95, 3.9, 1.25, size=8, color=GREEN)

callout(s, "NUDGE: Telegram üzerinden eksik PTI şoförlerine uyarı gönderir", 0.15, 3.55, 8.8, AMBER)
callout(s, "EXPORT CSV: Tüm yük listesini indirip muhasebe'ye gönder", 0.15, 3.95, 8.8, INK_DIM)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — PTI Inspection
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
topbar(s, "PTI Inspection")

panel_head(s, "PTI INSPECTION · BUGÜN · 28 MAYIS 2026", 0, 0.76, 9.5, "1/1 SUBMITTED")

# PTI row
add_rect(s, 0, 1.04, 9.5, 0.28, fill=PANEL, line_color=LINE)
for j, h in enumerate(["Şoför","Araç","Durum","Gönderim Saati","Fotoğraf","Doğrulama"]):
    add_text(s, h, 0.15 + j*1.55, 1.07, 1.45, 0.22, size=7, color=INK_MUTE)

add_rect(s, 0, 1.32, 9.5, 0.5, fill=BG, line_color=LINE)
row = [("Mustafa Bayram",INK),("TIR 9852",AMBER),("PTI · OK",GREEN),
       ("07:42 CT",INK_DIM),("6 FOTOĞRAF",INK_DIM),("✓ VALID",GREEN)]
for j, (v, c) in enumerate(row):
    add_text(s, v, 0.15 + j*1.55, 1.38, 1.45, 0.35, size=9, color=c)

# Fotoğraf slotları
add_text(s, "· FOTOĞRAF KANITI", 0.15, 1.95, 3, 0.25, size=8, color=INK_MUTE)
for i in range(6):
    add_rect(s, 0.15 + i*1.55, 2.2, 1.4, 1.0, fill=PANEL, line_color=LINE)
    labels = ["ÖN SOL","ÖN SAĞ","MOTOR","FREN","KARGO","PLAKA"]
    add_text(s, labels[i], 0.15 + i*1.55 + 0.1, 2.62, 1.2, 0.2,
             size=7, color=INK_MUTE, align=PP_ALIGN.CENTER)

# Sağ panel — açıklamalar
add_rect(s, 9.6, 0.76, 3.6, 6.5, fill=PANEL, line_color=LINE)
add_text(s, "· PTI NASIL ÇALIŞIR?", 9.75, 0.9, 3.3, 0.25, size=9, bold=True, color=AMBER)
steps = [
    ("Şoför Telegram'ı açar", "Her sabah /pti komutunu gönderir"),
    ("Fotoğraf çeker", "Min. 4 fotoğraf: ön, motor, fren, kargo"),
    ("Bot alır & kaydeder", "Sistem fotoğrafları otomatik loglar"),
    ("Dashboard güncellenir", "PTI · OK rozeti anında görünür"),
    ("NUDGE butonu", "Eksik PTI'lar için Telegram uyarısı gönderir"),
    ("Doğrulama", "Dispatcher VALID / INVALID işaretler"),
]
for i, (t, d) in enumerate(steps):
    y = 1.25 + i*0.85
    add_text(s, f"  {i+1}.", 9.75, y, 0.35, 0.3, size=10, color=AMBER, bold=True)
    add_text(s, t, 10.1, y, 3.0, 0.28, size=9, bold=True, color=INK)
    add_text(s, d, 10.1, y + 0.3, 3.0, 0.3, size=8, color=INK_DIM)
    if i < len(steps)-1:
        add_rect(s, 9.75, y + 0.65, 3.3, 0.01, fill=LINE, line_color=LINE)

callout(s, "Eksik PTI olan şoför o gün yola çıkamaz — compliance zorunludur", 0.15, 3.35, 9.3, RED)
callout(s, "Fotoğraflar sunucuda saklanır, istenildiğinde PDF rapor alınabilir", 0.15, 3.75, 9.3, INK_DIM)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — Loads / Dispatch
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
topbar(s, "Loads / Dispatch")

panel_head(s, "LOADS / DISPATCH", 0, 0.76, 8.5, "4 KAYIT")

# Filter bar
add_rect(s, 0, 1.04, 8.5, 0.35, fill=PANEL, line_color=LINE)
for i, (lbl, val, c) in enumerate([
    ("ALL","●",AMBER), ("PENDING","1",INK), ("IN TRANSIT","1",AMBER),
    ("DELIVERED","2",GREEN), ("CANCELLED","0",INK_MUTE)
]):
    add_rect(s, 0.05 + i*1.5, 1.08, 1.4, 0.25, fill=BG if i>0 else AMBER, line_color=LINE)
    add_text(s, f"{lbl}  {val}", 0.1 + i*1.5, 1.11, 1.3, 0.18, size=8,
             color=BG if i==0 else INK, align=PP_ALIGN.CENTER)

# Search input
add_rect(s, 7.55, 1.08, 0.9, 0.25, fill=BG, line_color=AMBER)
add_text(s, "🔍 ARA", 7.6, 1.11, 0.8, 0.18, size=7, color=AMBER)

# Tablo
add_rect(s, 0, 1.39, 8.5, 0.28, fill=PANEL, line_color=LINE)
headers = ["Yük #", "Araç", "Broker", "Güzergah", "Tarih", "Rate", "Durum", "Belgeler"]
for j, h in enumerate(headers):
    add_text(s, h, 0.1 + j*1.05, 1.42, 0.95, 0.22, size=7, color=INK_MUTE)

rows = [
    ("L-004","9852","XPO","Dallas TX→Chicago IL","2026-05-15","$2,400","IN TRANSIT",AMBER),
    ("L-003","9852","Echo","Dallas TX→Chicago IL","2026-04-20","$2,100","DELIVERED",GREEN),
]
for i, (ln, truck, broker, lane, date, rate, status, sc) in enumerate(rows):
    y = 1.67 + i*0.5
    add_rect(s, 0, y, 8.5, 0.48, fill=BG, line_color=LINE)
    vals = [(ln,AMBER),(truck,INK_DIM),(broker,INK),(lane,INK),(date,INK_DIM),(rate,GREEN)]
    for j, (v, c) in enumerate(vals):
        add_text(s, v, 0.1 + j*1.05, y+0.12, 0.95, 0.28, size=8, color=c)
    add_tag(s, status, 0.1 + 6*1.05, y+0.12, color=sc, size=7)
    add_text(s, "📎 2", 0.1 + 7*1.05, y+0.12, 0.5, 0.28, size=8, color=INK_DIM)

# Butonlar
btns = [("+ LOAD EKLE",AMBER,0.15,2.82),("▸ PARSE RC",GREEN,1.85,2.82),
        ("DAT BOARD",INK_DIM,3.55,2.82)]
for lbl, c, bx, by in btns:
    add_rect(s, bx, by, 1.55, 0.3, fill=c if c==AMBER else BG, line_color=c)
    add_text(s, lbl, bx+0.08, by+0.05, 1.4, 0.22, size=8,
             color=BG if c==AMBER else c, align=PP_ALIGN.CENTER)

# PARSE RC açıklama
add_rect(s, 0, 3.2, 8.5, 1.8, fill=PANEL, line_color=LINE)
add_text(s, "PARSE RC — Rate Confirmation PDF Otomatik Tanıma",
         0.15, 3.28, 6, 0.28, size=10, bold=True, color=GREEN)
add_text(s,
    "Broker'dan gelen PDF rate confirmation'ı sisteme sürükle-bırak ya da yükle.\n"
    "Sistem otomatik olarak şu alanları çıkarır:\n"
    "  ·  Yük numarası (load #)      ·  Broker adı\n"
    "  ·  Origin / Destination        ·  Pickup & Delivery tarihleri\n"
    "  ·  Rate ($)                         ·  Mesafe (miles)\n"
    "Alanlar düzenlenebilir formda gelir → LOAD OLUŞTUR butonuyla kayıt açılır.",
    0.15, 3.6, 8.2, 1.3, size=9, color=INK_DIM)

# Sağ panel
add_rect(s, 8.6, 0.76, 4.6, 6.5, fill=PANEL, line_color=LINE)
add_text(s, "· DAT LOAD BOARD", 8.75, 0.9, 4.3, 0.25, size=9, bold=True, color=AMBER)
add_text(s, "Mock mod (credentials bekleniyor)", 8.75, 1.18, 4.3, 0.22, size=8, color=AMBER)

dat_items = [
    ("Origin → Dest gir", "Arama alanına şehir yaz"),
    ("+ IMPORT butonu", "DAT yükünü direkt sisteme çeker"),
    ("Rate Intelligence", "Piyasa rayiç fiyatını gösterir"),
    ("Equipment filter", "Van / Reefer / Flatbed seçimi"),
]
for i, (t, d) in enumerate(dat_items):
    y = 1.5 + i*0.75
    add_text(s, "▸ " + t, 8.75, y, 4.2, 0.28, size=9, bold=True, color=AMBER)
    add_text(s, d, 8.75, y+0.28, 4.2, 0.28, size=8, color=INK_DIM)

add_text(s, "DAT BOARD NE İŞE YARAR?", 8.75, 4.6, 4.2, 0.25, size=9, bold=True, color=AMBER)
add_text(s,
    "Broker aramaya gerek kalmadan doğrudan\n"
    "DAT platformundaki yük ilanlarını gösterir.\n"
    "Piyasa fiyatıyla karşılaştırarak pazarlık\n"
    "gücü sağlar. Credentials gelince mock→live.",
    8.75, 4.88, 4.2, 1.2, size=8.5, color=INK_DIM)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — Fleet Roster
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
topbar(s, "Fleet Roster")

# Trucks panel
panel_head(s, "TRUCK ROSTER · 1 UNITS", 0, 0.76, 6.5, "+ ADD UNIT")
add_rect(s, 0, 1.04, 6.5, 0.28, fill=PANEL, line_color=LINE)
for j, h in enumerate(["Unit","Plate","Make/Model","YR","Odometer","Status","Driver"]):
    add_text(s, h, 0.1 + j*0.9, 1.07, 0.85, 0.22, size=7, color=INK_MUTE)
add_rect(s, 0, 1.32, 6.5, 0.72, fill=BG, line_color=LINE)
truck_row = [("9852",AMBER),("ABC9852",INK_DIM),("Volvo VNL860 '20",INK),
             ("2020",INK),("125,400 MI",INK_DIM),("ACTIVE",GREEN),("M. Bayram",INK)]
for j, (v, c) in enumerate(truck_row):
    add_text(s, v, 0.1 + j*0.9, 1.42, 0.85, 0.4, size=8, color=c)
# Fuel bar
add_rect(s, 0.1, 1.72, 0.6, 0.06, fill=LINE, line_color=LINE)
add_rect(s, 0.1, 1.72, 0.48, 0.06, fill=GREEN, line_color=GREEN)
add_text(s, "80%", 0.72, 1.7, 0.3, 0.14, size=7, color=INK_MUTE)

callout(s, "Satıra tıkla → Düzenleme modalı: status, yakıt, kilometre, VIN güncelle", 0.1, 2.15, 6.3, AMBER)

# Drivers panel
panel_head(s, "DRIVER ROSTER · 1 ACTIVE", 6.6, 0.76, 3.0, "+ ADD DRIVER")
add_rect(s, 6.6, 1.04, 3.0, 1.2, fill=BG, line_color=LINE)
add_rect(s, 6.7, 1.12, 0.45, 0.45, fill=PANEL, line_color=LINE)
add_text(s, "MB", 6.72, 1.18, 0.4, 0.35, size=13, bold=True, color=AMBER,
         font_name="Arial Black", align=PP_ALIGN.CENTER)
add_text(s, "Mustafa Bayram", 7.2, 1.15, 2.3, 0.25, size=9, color=INK)
add_text(s, "TG · 6290812296  ·  TIR 9852", 7.2, 1.4, 2.3, 0.22, size=8, color=INK_DIM)
add_text(s, "CDL · B123456  ·  EXP 2027-01", 7.2, 1.62, 2.3, 0.22, size=8, color=GREEN)
callout(s, "CDL expiry kırmızıya döner 90 gün önce uyarır", 6.6, 2.35, 6.7, AMBER)

# Maintenance panel
panel_head(s, "MAINTENANCE LOG", 9.7, 0.76, 3.5, "+ LOG SERVICE")
maint = [
    ("OIL","9852","2026-04-15","Valvoline 15W40","$180","COMPLETED",GREEN),
    ("TIRE","9852","2026-03-10","Michelin X","$850","COMPLETED",GREEN),
    ("INSP","9852","2026-02-01","Annual DOT","$120","COMPLETED",GREEN),
]
for i, (typ, truck, date, desc, cost, status, sc) in enumerate(maint):
    y = 1.04 + i*0.55
    add_rect(s, 9.7, y, 3.5, 0.52, fill=BG, line_color=LINE)
    add_tag(s, typ, 9.78, y+0.1, color=AMBER, size=7)
    add_text(s, truck, 10.25, y+0.08, 0.5, 0.2, size=9, color=AMBER)
    add_text(s, desc, 10.78, y+0.08, 2.3, 0.2, size=8, color=INK)
    add_text(s, date, 9.78, y+0.3, 1.2, 0.2, size=7, color=INK_MUTE)
    add_text(s, cost, 11.5, y+0.08, 0.6, 0.2, size=8, color=GREEN, align=PP_ALIGN.RIGHT)

# Edit modal örneği
add_rect(s, 0.15, 2.55, 6.3, 4.6, fill=PANEL, line_color=LINE)
add_text(s, "· DÜZENLEME MODALI — TIR 9852", 0.3, 2.68, 5, 0.28, size=9, bold=True, color=INK_DIM)
fields = [("Plate","ABC9852"),("Make","Volvo"),("Model","VNL860"),("VIN","1FUJ..."),
          ("Status","ACTIVE ▾"),("Fuel %","80"),("Odometer","125400")]
for i, (k, v) in enumerate(fields):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 3.1
    y = 3.05 + row * 0.72
    add_text(s, k, x, y, 1.2, 0.2, size=7, color=INK_MUTE)
    add_rect(s, x, y+0.22, 2.85, 0.32, fill=BG, line_color=LINE)
    add_text(s, v, x+0.08, y+0.27, 2.7, 0.22, size=9, color=INK)
add_rect(s, 0.3, 5.65, 5.9, 0.35, fill=AMBER, line_color=AMBER)
add_text(s, "▸ SAVE CHANGES", 0.3, 5.7, 5.9, 0.25, size=10, bold=True,
         color=BG, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — Factoring
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
topbar(s, "Factoring")

# Summary stats
for i, (lbl, val, sub, c) in enumerate([
    ("Total Invoices","02","kayıt",INK),
    ("Pending","$0","gönderilmedi",AMBER),
    ("Submitted","$0","onay bekliyor",INK_DIM),
    ("Paid","$225","ödendi",GREEN),
    ("Fee Total","$8","RTS komisyonu",INK_MUTE),
]):
    x = i * 2.6
    add_rect(s, x, 0.76, 2.55, 0.8, fill=PANEL, line_color=LINE)
    add_text(s, lbl, x+0.12, 0.82, 2.3, 0.2, size=8, color=INK_MUTE)
    add_text(s, val, x+0.12, 1.0, 2.3, 0.38, size=20, bold=True, color=c,
             font_name="Arial Black")
    add_text(s, sub, x+0.12, 1.4, 2.3, 0.18, size=7, color=INK_MUTE)

panel_head(s, "FACTORING KAYITLARI", 0, 1.56, 9.5, "· RTS FINANCIAL")

# Tablo
add_rect(s, 0, 1.84, 9.5, 0.28, fill=PANEL, line_color=LINE)
for j, h in enumerate(["Yük #","Broker","Güzergah","Rate","Invoice","Net","RTS Status","Ödeme"]):
    add_text(s, h, 0.1 + j*1.18, 1.87, 1.1, 0.22, size=7, color=INK_MUTE)

fact_rows = [
    ("L-003","Echo","DAL→CHI","$2,100","$2,100","$2,026","PAID","2026-05-01",GREEN),
    ("L-004","XPO","DAL→CHI","$2,400","—","—","PENDING","—",AMBER),
]
for i, (*vals, sc) in enumerate(fact_rows):
    y = 2.12 + i*0.52
    add_rect(s, 0, y, 9.5, 0.5, fill=BG, line_color=LINE)
    for j, v in enumerate(vals):
        add_text(s, v, 0.1+j*1.18, y+0.12, 1.1, 0.3, size=8,
                 color=sc if j in [6,7] else INK if j<4 else GREEN)

# RTS iş akışı
add_rect(s, 0, 3.25, 9.5, 3.95, fill=PANEL, line_color=LINE)
add_text(s, "RTS FACTORING İŞ AKIŞI — ADIM ADIM",
         0.15, 3.35, 7, 0.3, size=11, bold=True, color=AMBER)

workflow = [
    ("Teslim Edilen Yük", "Load status → DELIVERED olur", AMBER),
    ("Invoice Oluştur", "Factoring sekmesinden 'Invoice Oluştur' → yük seç, tutar gir", GREEN),
    ("RTS'e Gönder", "'Submit to RTS' → status: PENDING → APPROVED değişir", AMBER),
    ("Ödeme Takibi", "RTS status PAID olunca net tutar görünür (%3.5 komisyon düşülmüş)", GREEN),
]
for i, (title, desc, c) in enumerate(workflow):
    x = 0.2 + i * 2.35
    y = 3.75
    add_rect(s, x, y, 2.2, 1.0, fill=BG, line_color=c)
    add_text(s, f"ADIM {i+1}", x+0.1, y+0.08, 2.0, 0.2, size=7, color=c)
    add_text(s, title, x+0.1, y+0.3, 2.0, 0.28, size=9, bold=True, color=INK)
    add_text(s, desc, x+0.1, y+0.6, 2.0, 0.38, size=7.5, color=INK_DIM)
    if i < 3:
        add_text(s, "→", x+2.22, y+0.45, 0.12, 0.3, size=14, color=AMBER)

add_text(s, "NET FORMÜL:", 0.2, 4.9, 2, 0.25, size=8, color=INK_MUTE)
add_text(s, "Net = Invoice × (1 – %3.5)  →  $2,100 × 0.965 = $2,026.50",
         0.2, 5.15, 9.0, 0.35, size=11, color=GREEN, font_name="Courier New")

# Sağ panel
add_rect(s, 9.6, 1.56, 3.6, 5.6, fill=PANEL, line_color=LINE)
add_text(s, "· FACTORING NEDEN?", 9.75, 1.7, 3.3, 0.25, size=9, bold=True, color=AMBER)
reasons = [
    "Teslimden 24 saat içinde nakit",
    "Broker 60 gün ödeme yaparsa\nsiz %3.5 karşılığı hemen alırsınız",
    "RTS Financial entegrasyonu",
    "Invoice numarası otomatik",
    "Tüm ödeme geçmişi kayıtlı",
]
for i, r in enumerate(reasons):
    add_text(s, "▸ " + r, 9.75, 2.05 + i*0.8, 3.3, 0.65, size=9, color=INK_DIM)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — IFTA
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
topbar(s, "IFTA")

add_rect(s, 0, 0.76, 13.33, 0.45, fill=PANEL, line_color=LINE)
add_text(s, "· IFTA MILEAGE SUMMARY ·", 0.2, 0.84, 3.5, 0.28, size=8, color=INK_MUTE)
for i, q in enumerate([1,2,3,4]):
    clr = AMBER if q==2 else BG
    add_rect(s, 4.0 + i*1.1, 0.84, 1.0, 0.28,
             fill=clr, line_color=AMBER if q==2 else LINE)
    add_text(s, f"Q{q}", 4.05 + i*1.1, 0.87, 0.9, 0.22, size=9,
             color=BG if q==2 else INK, align=PP_ALIGN.CENTER)
add_text(s, "TOTAL MILES  6,172", 8.0, 0.84, 2.5, 0.28, size=9, color=AMBER)
add_text(s, "EST. TAX  $1,240", 10.6, 0.84, 2.5, 0.28, size=9, color=RED)

# State tablo
panel_head(s, "STATE MILEAGE BREAKDOWN · Q2 2026", 0, 1.21, 9.5, "6 STATES")
add_rect(s, 0, 1.49, 9.5, 0.28, fill=PANEL, line_color=LINE)
for j, h in enumerate(["State","Miles","% of Total","Tax Rate","Gallons","Est. Tax"]):
    add_text(s, h, 0.15 + j*1.55, 1.52, 1.45, 0.22, size=7, color=INK_MUTE)

state_data = [
    ("TX", 920, "29.8%", "$0.200", 141, "$28"),
    ("IL", 920, "29.8%", "$0.462", 141, "$65"),
    ("WI", 500, "16.2%", "$0.329", 77, "$25"),
    ("NY", 500, "16.2%", "$0.451", 77, "$35"),
]
for i, (st, mi, pct, rate, gal, tax) in enumerate(state_data):
    y = 1.77 + i*0.5
    add_rect(s, 0, y, 9.5, 0.48, fill=BG, line_color=LINE)
    add_tag(s, st, 0.1, y+0.12, color=AMBER, size=8)
    bar_w = mi / 1000.0 * 1.2
    add_rect(s, 1.7, y+0.2, 1.2, 0.08, fill=LINE, line_color=LINE)
    add_rect(s, 1.7, y+0.2, bar_w, 0.08, fill=AMBER, line_color=AMBER)
    add_text(s, str(mi), 2.95, y+0.12, 0.5, 0.28, size=9, color=INK)
    for j, (v, c) in enumerate([(pct,INK_MUTE),(rate,INK),(str(gal),INK_DIM),(tax,RED)]):
        add_text(s, v, 3.55 + j*1.55, y+0.12, 1.45, 0.28, size=9, color=c)

# Sağ panel
add_rect(s, 9.6, 1.21, 3.6, 6.05, fill=PANEL, line_color=LINE)
add_text(s, "· IFTA NEDİR?", 9.75, 1.35, 3.3, 0.25, size=9, bold=True, color=AMBER)
add_text(s,
    "International Fuel Tax Agreement.\n\n"
    "Her eyalette kaç mil sürdüğünü bildirip\n"
    "o eyaletin yakıt vergisini ödeme sistemi.\n\n"
    "Quarterly (3 ayda bir) dosyalama zorunlu.",
    9.75, 1.65, 3.3, 1.5, size=9, color=INK_DIM)

deadlines = [("Q1 Jan–Mar","Apr 30"),("Q2 Apr–Jun","Jul 31"),
             ("Q3 Jul–Sep","Oct 31"),("Q4 Oct–Dec","Jan 31")]
add_text(s, "· TERMİNLER", 9.75, 3.25, 3.3, 0.25, size=9, bold=True, color=AMBER)
for i, (q, d) in enumerate(deadlines):
    y = 3.55 + i * 0.52
    add_text(s, q, 9.75, y, 2.0, 0.25, size=9, color=INK_DIM)
    add_text(s, d, 11.8, y, 1.3, 0.25, size=9, color=AMBER, align=PP_ALIGN.RIGHT)
    add_rect(s, 9.75, y+0.28, 3.3, 0.01, fill=LINE, line_color=LINE)

add_rect(s, 9.75, 5.6, 3.3, 1.4, fill=BG, line_color=AMBER)
add_text(s, "⚠ UYARI", 9.85, 5.68, 3.1, 0.22, size=8, bold=True, color=AMBER)
add_text(s,
    "Mil tahmini origin→dest 50/50\nbölüşüm ile hesaplanır. Gerçek\nIFTA için GPS log şart.",
    9.85, 5.92, 3.1, 0.9, size=8, color=AMBER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — Multi-Tenant Mimari
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
add_rect(s, 0, 0, 0.06, 7.5, fill=AMBER, line_color=AMBER)

add_text(s, "MULTI-TENANT MİMARİ", 0.3, 0.2, 10, 0.5,
         size=22, bold=True, color=AMBER, font_name="Arial Black")
add_text(s, "Birden fazla filo şirketi — tek platform, tam izolasyon",
         0.3, 0.72, 10, 0.3, size=11, color=INK_DIM)
add_rect(s, 0.3, 1.05, 12.7, 0.01, fill=LINE, line_color=LINE)

# Sol — Şirket A
add_rect(s, 0.3, 1.2, 5.8, 5.4, fill=PANEL, line_color=GREEN)
add_text(s, "ŞİRKET A — TIR FLEET", 0.5, 1.35, 5.4, 0.3, size=11, bold=True, color=GREEN)
a_items = ["Tır: 9852, 9853, 9854","Şoför: M.Bayram, A.Demir","Yük: L-001 → L-045","Factoring: RTS bağlı"]
for i, item in enumerate(a_items):
    add_text(s, "  ▸ " + item, 0.5, 1.8 + i*0.45, 5.4, 0.35, size=10, color=INK_DIM)
add_text(s, "company_id = 1", 0.5, 3.8, 5.4, 0.3, size=9,
         color=GREEN, font_name="Courier New")
add_text(s, "admin@tirfleet.com  |  OWNER rolü", 0.5, 4.1, 5.4, 0.3, size=9, color=INK_MUTE)

# Orta — Platform
add_rect(s, 6.2, 2.0, 0.9, 3.0, fill=BG, line_color=LINE)
add_text(s, "F\nL\nE\nE\nT\nS\nY\nN\nC", 6.25, 2.1, 0.8, 2.8, size=9,
         color=AMBER, bold=True, align=PP_ALIGN.CENTER, font_name="Arial Black")

# Sağ — Şirket B
add_rect(s, 7.2, 1.2, 5.8, 5.4, fill=PANEL, line_color=AMBER)
add_text(s, "ŞİRKET B — YILMAZ NAKLİYE", 7.4, 1.35, 5.4, 0.3, size=11, bold=True, color=AMBER)
b_items = ["Tır: Y-001, Y-002","Şoför: K.Yılmaz","Yük: L-001 → L-012","Factoring: ayrı"]
for i, item in enumerate(b_items):
    add_text(s, "  ▸ " + item, 7.4, 1.8 + i*0.45, 5.4, 0.35, size=10, color=INK_DIM)
add_text(s, "company_id = 2", 7.4, 3.8, 5.4, 0.3, size=9,
         color=AMBER, font_name="Courier New")
add_text(s, "info@yilmaznakliye.com  |  OWNER rolü", 7.4, 4.1, 5.4, 0.3, size=9, color=INK_MUTE)

add_text(s, "✗  Şirket B, Şirket A'nın HİÇBİR verisini göremez",
         0.3, 6.75, 12.5, 0.35, size=11, color=RED, align=PP_ALIGN.CENTER)

add_text(s, "SaaS modeli: Her yeni filo /auth/register ile şirket kurar, kendi izole ortamında çalışır.",
         0.3, 7.1, 12.5, 0.28, size=9, color=INK_MUTE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — Teknik Altyapı
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
add_rect(s, 0, 0, 0.06, 7.5, fill=AMBER, line_color=AMBER)

add_text(s, "TEKNİK ALTYAPI", 0.3, 0.2, 10, 0.5,
         size=22, bold=True, color=AMBER, font_name="Arial Black")
add_rect(s, 0.3, 0.72, 12.7, 0.01, fill=LINE, line_color=LINE)

tech_blocks = [
    ("BACKEND", "FastAPI (Python)",
     ["SQLAlchemy ORM","SQLite veritabanı","JWT authentication","bcrypt şifreleme","pdfplumber (PDF parse)"]),
    ("FRONTEND", "React + Vite",
     ["TanStack Query (data fetch)","React Router v6","Dark terminal UI","Auto-refresh (30s)","Responsive CSS"]),
    ("BOT", "Telegram Bot",
     ["python-telegram-bot","PTI fotoğraf alımı","Şoför nudge","Broadcast mesaj","Auto pti logging"]),
    ("DEPLOY", "Hetzner VPS",
     ["Nginx reverse proxy","systemd servisler","SSL ready","GitHub CI/CD","5.161.61.221"]),
]
for i, (title, sub, items) in enumerate(tech_blocks):
    col = i % 2
    row = i // 2
    l = 0.3 + col * 6.5
    t = 0.9 + row * 3.1

    add_rect(s, l, t, 6.2, 2.9, fill=PANEL, line_color=LINE)
    add_rect(s, l, t, 6.2, 0.52, fill=LINE, line_color=LINE)
    add_text(s, title, l+0.15, t+0.08, 3, 0.32, size=13, bold=True, color=AMBER,
             font_name="Arial Black")
    add_text(s, sub, l+0.15, t+0.3, 4, 0.2, size=8, color=INK_DIM)
    for j, item in enumerate(items):
        add_text(s, "▸ " + item, l+0.2, t + 0.65 + j*0.42, 5.7, 0.35, size=9, color=INK)


# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — Kapanış / Özet
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
add_rect(s, 0, 0, 0.06, 7.5, fill=AMBER, line_color=AMBER)

add_text(s, "FLEETSYNC", 0.4, 1.0, 12, 1.2,
         size=64, bold=True, color=AMBER, font_name="Arial Black")
add_text(s, "Production-ready. Canlıda çalışıyor.", 0.45, 2.2, 10, 0.45,
         size=16, color=INK)

features = [
    ("✓ PTI","Telegram bot ile sabah check-in"),
    ("✓ Dispatch","Yük yönetimi, DAT, PDF parse"),
    ("✓ Fleet","Tır + şoför + bakım + CDL takibi"),
    ("✓ Factoring","RTS invoice iş akışı"),
    ("✓ IFTA","Quarterly mileage raporu"),
    ("✓ Multi-Tenant","Sınırsız şirket, tam izolasyon"),
    ("✓ Auth","JWT, rol bazlı erişim"),
    ("✓ Mobile","Responsive, her cihazdan erişim"),
]
for i, (feat, desc) in enumerate(features):
    col = i % 2
    row = i // 2
    x = 0.45 + col * 6.3
    y = 2.85 + row * 0.72
    add_text(s, feat, x, y, 1.6, 0.3, size=11, bold=True, color=GREEN)
    add_text(s, desc, x + 1.65, y, 4.5, 0.3, size=10, color=INK_DIM)

add_rect(s, 0.45, 5.9, 12.4, 0.01, fill=LINE, line_color=LINE)
add_text(s, "http://5.161.61.221", 0.45, 6.05, 5, 0.35, size=12, color=AMBER,
         font_name="Courier New")
add_text(s, "contact@fleetsync.io", 7.0, 6.05, 5.8, 0.35, size=12, color=INK_DIM,
         align=PP_ALIGN.RIGHT)
add_text(s, "© 2026 FLEETSYNC · TÜM HAKLARI SAKLIDIR", 0.45, 6.5, 12, 0.3,
         size=8, color=INK_MUTE, align=PP_ALIGN.CENTER)


# ── Kaydet ──────────────────────────────────────────────────────
out = r"C:\Users\berki\Desktop\FLEETSYNC_SUNUM.pptx"
prs.save(out)
print("Sunum kaydedildi:", out)
print("Slide sayisi:", len(prs.slides))
